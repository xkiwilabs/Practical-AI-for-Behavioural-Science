"""
Build Week 1 lecture slides using the MQ PowerPoint template.

Usage:
    python generate_figures.py   # run first to create images
    python build_slides.py

Output:
    ../slides/Week01_From_Mind_to_Model.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Paths
TEMPLATE_PATH = "../../../.old_materials/Decoding Human Behaviour.pptx"
OUTPUT_DIR = "../slides"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Week01_From_Mind_to_Model.pptx")
IMAGES_DIR = "images"

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Layout indices (from template analysis)
LAYOUT_TITLE = 0        # Title Slide
LAYOUT_SECTION = 3      # Section Header
LAYOUT_CONTENT = 4      # One Column Content (title + content + footer)
LAYOUT_TITLE_ONLY = 8   # Title Only
LAYOUT_PICTURE = 9      # Picture layout
LAYOUT_BLANK = 10       # Blank

# Colours
MQ_RED = RGBColor(0xA7, 0x19, 0x30)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREY_TEXT = RGBColor(0x8C, 0x8C, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x4A, 0x90, 0xD9)


def load_template():
    """Load the MQ template and strip all existing slides."""
    prs = Presentation(TEMPLATE_PATH)
    # Remove all existing slides by deleting their XML references.
    # We keep the slide masters and layouts intact.
    while len(prs.slides._sldIdLst) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId is None:
            rId = sldId.get("r:id")
        prs.slides._sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except (KeyError, ValueError):
                pass
    return prs


def add_title_slide(prs, title, subtitle, footer=""):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 15:  # Subtitle area
            ph.text = subtitle
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
        elif ph.placeholder_format.idx == 14:  # Top text
            ph.text = footer
    return slide


def add_section_slide(prs, title, subtitle=""):
    """Add a section header slide (layout 3)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 14 and subtitle:
            ph.text = subtitle
    return slide


def add_content_slide(prs, title, bullets, subtitle_text="", image_path=None):
    """Add a content slide with title and bullet points.
    If image_path given, places image on right side."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(20)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 15 and subtitle_text:
            ph.text = subtitle_text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = MQ_RED
        elif ph.placeholder_format.idx == 1:  # Content area
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()

                # Support bold prefix with "**text**rest" syntax
                if bullet.startswith("**") and "**" in bullet[2:]:
                    end = bullet.index("**", 2)
                    bold_part = bullet[2:end]
                    rest = bullet[end + 2:]
                    run_bold = para.add_run()
                    run_bold.text = bold_part
                    run_bold.font.bold = True
                    run_bold.font.size = Pt(18)
                    if rest:
                        run_rest = para.add_run()
                        run_rest.text = rest
                        run_rest.font.size = Pt(18)
                else:
                    run = para.add_run()
                    run.text = bullet
                    run.font.size = Pt(18)

                para.space_after = Pt(6)

    if image_path and os.path.exists(image_path):
        # Place image on right half of slide
        slide.shapes.add_picture(
            image_path,
            left=Inches(6.5), top=Inches(1.8),
            width=Inches(5.5)
        )

    return slide


def add_image_slide(prs, title, image_path, caption="", subtitle_text=""):
    """Add a slide with a large centred image."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(20)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 15 and subtitle_text:
            ph.text = subtitle_text

    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            left=Inches(1.0), top=Inches(1.5),
            width=Inches(11.0)
        )

    if caption:
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(6.8), Inches(11.0), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = GREY_TEXT

    return slide


def add_discussion_slide(prs, question, context=""):
    """Add a 'Think About It' discussion slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = question
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 14:
            ph.text = "Think about it..."
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
                    run.font.italic = True

    return slide


def add_two_column_text_slide(prs, title, left_title, left_bullets,
                               right_title, right_bullets, subtitle_text=""):
    """Add a two-column text slide using Title Only layout + textboxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(20)
                    run.font.bold = True
        elif ph.placeholder_format.idx == 15 and subtitle_text:
            ph.text = subtitle_text

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = left_title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = BLUE
    for bullet in left_bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(15)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = right_title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = MQ_RED
    for bullet in right_bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(15)
        p.space_after = Pt(4)

    return slide


def img(name):
    """Helper to get image path."""
    return os.path.join(IMAGES_DIR, name)


def build_deck():
    prs = load_template()

    # ============================================================
    # PART 1: COURSE OVERVIEW (~12 slides)
    # ============================================================

    add_title_slide(
        prs,
        "Week 1 — From Mind to Model",
        "Why ML Belongs in Psychological Science",
        "PSYC4411 : CURRENT ADVANCES IN PSYCHOLOGICAL METHODS AND ANALYSES"
    )

    add_content_slide(prs, "Welcome to PSYC4411", [
        "This course: modern ML/AI tools for psychological research",
        "No coding or technical background needed",
        "You'll use AI assistants to help write code",
        "Goal: become a researcher who uses these tools thoughtfully",
    ], subtitle_text="WHAT THIS COURSE IS ABOUT")

    add_content_slide(prs, "What You'll Learn", [
        "Prediction — regression, classification models",
        "Structure — clustering, dimensionality reduction",
        "Neural networks — when complexity is warranted",
        "Text & language — embeddings and LLMs",
        "Critical thinking — what ML can and can't tell us",
    ], subtitle_text="COURSE TOPICS")

    add_content_slide(prs, "Course Structure", [
        "13 weeks: alternating lectures and labs",
        "Lecture weeks (odd): concepts, theory, discussion",
        "Lab weeks (even): hands-on group challenges with AI assistants",
        "Week 11: hybrid lecture + lab",
        "Weeks 12–13: viva review and final discussion",
    ], subtitle_text="HOW THE COURSE WORKS")

    add_content_slide(prs, "Assessment", [
        "**Written assignment (100%)** — due Week 8",
        "**Viva exam** — Weeks 12–13",
        "No traditional coding exam",
        "Group challenge presentations each lab week (formative)",
    ], subtitle_text="HOW YOU'RE ASSESSED")

    add_content_slide(prs, "The Tools You'll Use", [
        "**Python** — the programming language (you won't memorise syntax)",
        "**VS Code** — your code editor",
        "**Jupyter Notebooks** — interactive documents mixing code and text",
        "**LLM Assistants** — ChatGPT, Claude, Copilot (your coding partners)",
        "**Research Tools** — NotebookLM, Elicit, Semantic Scholar",
    ], subtitle_text="YOUR TOOLKIT")

    add_content_slide(prs, "Homework Before Week 2", [
        "Complete the Getting Started setup guide",
        "Create a GitHub account (personal email, not MQ)",
        "Apply for the Student Developer Pack (free Copilot Pro)",
        "Install VS Code + Python + run the setup script",
        "Sign up for at least one AI assistant",
        "If you get stuck — bring your laptop to Week 2",
    ], subtitle_text="SETUP")

    # ============================================================
    # PART 2: WHAT ARE AI AND ML? (~10 slides)
    # ============================================================

    add_section_slide(prs, "What Are AI and\nMachine Learning?")

    add_content_slide(prs, "Definitions", [
        "**Artificial Intelligence (AI)** — systems that perform tasks "
        "requiring human-like intelligence",
        "**Machine Learning (ML)** — a subset of AI that learns patterns "
        "from data instead of following explicit rules",
        "**Deep Learning** — ML using neural networks (layers of "
        "mathematical operations)",
        "**Generative AI** — systems that create new content "
        "(text, images, code)",
    ], subtitle_text="KEY TERMS")

    add_image_slide(
        prs,
        "The AI Family Tree",
        img("ai_hierarchy.png"),
        subtitle_text="AI > ML > DEEP LEARNING > GENERATIVE AI"
    )

    add_content_slide(prs, "How Humans Learn", [
        "A child learns to recognise dogs from a few examples",
        "No rule book — pattern learning from experience",
        "ML does the same thing, but mathematically",
        "Key difference: humans need ~5 examples",
        "ML models typically need thousands or millions",
    ], subtitle_text="A PSYCHOLOGY ANALOGY")

    add_discussion_slide(
        prs,
        "Why can a toddler learn \"dog\" from a few examples\n"
        "while an ML model needs thousands?\n\n"
        "What does the human bring to the task\nthat the model doesn't?"
    )

    add_image_slide(
        prs,
        "Traditional Programming vs. Machine Learning",
        img("traditional_vs_ml.png"),
        subtitle_text="TWO APPROACHES TO COMPUTING"
    )

    # ============================================================
    # PART 3: BRIEF HISTORY (~6 slides)
    # ============================================================

    add_section_slide(prs, "A Brief History of AI",
                      "From Turing to Transformers")

    add_image_slide(
        prs,
        "The Story So Far",
        img("ai_timeline.png"),
    )

    add_content_slide(prs, "The Transformer Moment (2017)", [
        "\"Attention Is All You Need\" — Vaswani et al.",
        "New architecture for processing sequences of text",
        "Foundation of every modern LLM",
        "GPT, Claude, Gemini — all built on transformers",
    ], subtitle_text="THE BREAKTHROUGH")

    add_content_slide(prs, "2022–2026: The Current Wave", [
        "2022: ChatGPT launches — millions discover LLMs",
        "2023: GPT-4, Claude 2, multimodal models",
        "2024: AI coding assistants go mainstream",
        "2025: Vibe coding, Deep Research, agentic AI",
        "2026: AI agents that plan and execute multi-step tasks",
    ], subtitle_text="WHERE WE ARE NOW")

    # ============================================================
    # PART 4: AI TOOLS IN 2026 (~8 slides)
    # ============================================================

    add_section_slide(prs, "The AI Tools Landscape\nin 2026",
                      "What's available to you right now")

    add_content_slide(prs, "Conversational AI + Deep Research", [
        "ChatGPT, Claude, Gemini — evolved beyond basic chat",
        "**Deep Research mode:** AI autonomously searches, reads, "
        "and synthesises across many sources",
        "Like a research assistant who reads dozens of papers for you",
        "Useful for literature reviews, understanding new methods",
    ], subtitle_text="CHATBOTS THAT RESEARCH")

    add_content_slide(prs, "Coding Assistance + Vibe Coding", [
        "**GitHub Copilot** — AI code suggestions inside VS Code",
        "**Vibe coding** (Karpathy, Feb 2025) — describe what you want "
        "in English, AI writes the code",
        "This is how you'll code in this course",
        "25% of YC Winter 2025 startups: 95% AI-generated code",
    ], subtitle_text="AI-ASSISTED PROGRAMMING")

    add_content_slide(prs, "Research Tools", [
        "**NotebookLM** — upload papers, get AI summaries and mind maps",
        "**Elicit** — evidence synthesis from 200M+ papers",
        "**Consensus** — research questions answered from peer-reviewed lit",
        "**Semantic Scholar** — AI-powered paper discovery",
    ], subtitle_text="BUILT FOR ACADEMICS")

    add_content_slide(prs, "Tools Augment, Not Replace", [
        "A hammer doesn't build a house by itself",
        "A skilled carpenter with a hammer builds better houses faster",
        "Same relationship with AI tools:",
        "You bring the research question, domain knowledge, "
        "and critical judgement",
        "AI brings speed, pattern-finding, and tireless computation",
    ], subtitle_text="THE KEY MINDSET")

    # ============================================================
    # PART 5: PROMPT ENGINEERING (~6 slides)
    # ============================================================

    add_section_slide(prs, "Prompt Engineering\nand Context Engineering",
                      "The most practical skill you'll learn this semester")

    add_two_column_text_slide(
        prs,
        "Weak Prompt vs. Strong Prompt",
        "Weak Prompt",
        [
            "\"Make me a graph\"",
            "",
            "AI has to guess:",
            "  What data?",
            "  What type of graph?",
            "  What library?",
            "  What variables?",
            "",
            "Result: probably wrong",
        ],
        "Strong Prompt",
        [
            "\"Create a scatter plot of Depression",
            "  score vs Sleep hours from my",
            "  DataFrame called data, with points",
            "  coloured by Gender, using matplotlib.",
            "  Add a title, axis labels, and a legend.\"",
            "",
            "Result: almost certainly correct",
            "on the first try",
        ],
        subtitle_text="PROMPT ENGINEERING"
    )

    add_content_slide(prs, "Context Engineering", [
        "The AI doesn't know your data, your question, or your constraints",
        "The more relevant context you provide, the better the output",
        "Think of it as briefing a new research assistant:",
        "What libraries are you using? What does your data look like?",
        "What have you already tried? What didn't work?",
    ], subtitle_text="GIVE THE AI WHAT IT NEEDS")

    add_content_slide(prs, "Key Principles", [
        "**Be specific** — format, details, exact variable names",
        "**Give context** — your data, tools, and goal",
        "**Ask for explanations** — \"...and explain each line of code\"",
        "**State constraints** — \"Use only pandas and matplotlib\"",
    ], subtitle_text="WORKING WITH AI ASSISTANTS")

    add_discussion_slide(
        prs,
        "What does the need for prompt engineering\n"
        "tell us about the difference between\n"
        "human understanding and what AI does?\n\n"
        "What would AI need to not require\ndetailed prompts?"
    )

    # ============================================================
    # PART 6: LLM PROBLEM-SOLVING LOOP (~5 slides)
    # ============================================================

    add_section_slide(prs, "The LLM Problem-Solving Loop",
                      "Your core workflow all semester")

    add_image_slide(
        prs,
        "Two Nested Loops",
        img("llm_loop.png"),
    )

    add_content_slide(prs, "Outer Loop — Your Research Process", [
        "**1. PLAN** — What question are you answering? What output?",
        "**2. EXECUTE** — Use the inner loop to get AI-generated code",
        "**3. EVALUATE** — Does the result answer your question? Correct?",
        "**4. DOCUMENT** — Record what you did and learned",
    ], subtitle_text="THE BIG PICTURE")

    add_content_slide(prs, "Inner Loop — Working with the AI", [
        "**1. ENGINEER** — Craft your prompt with context",
        "**2. PROMPT** — Send it to the AI",
        "**3. VERIFY** — Run the code. Does it work? Look right?",
        "**4. REFINE** — Fix, add context, re-prompt",
        "",
        "This loop runs 2–5 times per task. That's normal.",
    ], subtitle_text="THE ITERATION CYCLE")

    add_content_slide(prs, "You Are Always in Control", [
        "You decide what to build",
        "You verify the output",
        "You judge whether it's correct",
        "The AI is powerful but doesn't understand your research "
        "question the way you do",
    ], subtitle_text="THE KEY INSIGHT")

    # ============================================================
    # PART 7: ML IN PSYCHOLOGICAL SCIENCE (~12 slides)
    # ============================================================

    add_section_slide(prs, "ML and AI in\nPsychological Science",
                      "Why should psychologists care?")

    add_image_slide(
        prs,
        "Prediction vs. Explanation",
        img("prediction_vs_explanation.png"),
    )

    add_content_slide(prs, "Prediction vs. Explanation", [
        "Psychology has focused on explanation — why does X cause Y?",
        "ML adds prediction — can we forecast Y from X?",
        "Yarkoni & Westfall (2017): both approaches strengthen each other",
        "This course bridges the two traditions",
    ], subtitle_text="TWO COMPLEMENTARY APPROACHES")

    add_discussion_slide(
        prs,
        "An ML model can predict depression risk\n"
        "from phone data, but it doesn't \"understand\"\n"
        "depression.\n\n"
        "A clinician might be less accurate but\n"
        "understands the patient.\n\n"
        "When is prediction without understanding\nenough? When is it not?"
    )

    add_content_slide(prs, "Example 1: Digital Phenotyping", [
        "Smartphone sensor data: GPS, screen time, sleep, typing speed",
        "ML models predict mental health episodes before they happen",
        "Passive data — no self-report needed",
        "Early intervention without waiting for a crisis",
    ], subtitle_text="MENTAL HEALTH")

    add_content_slide(prs, "Example 2: LLMs as Cognitive Models", [
        "Hagendorff, Fabi, & Kosinski (2023) — Nature Computational Science",
        "Tested LLMs on classic cognitive psychology tasks",
        "Larger models developed human-like reasoning biases",
        "But biases disappeared in ChatGPT (\"System 2\" reasoning)",
        "Same behaviour — but the same mechanism?",
    ], subtitle_text="COGNITIVE SCIENCE")

    add_discussion_slide(
        prs,
        "If an LLM shows the same reasoning bias\n"
        "as a human on a classic cognitive task,\n"
        "does that mean it's \"thinking\" the same way?\n\n"
        "What evidence would you need\nto make that claim?"
    )

    add_content_slide(prs, "Example 3: Silicon Samples", [
        "Argyle et al. (2023) — Political Analysis",
        "LLMs conditioned with demographic backstories",
        "Simulated human survey responses with surprising accuracy",
        "\"Algorithmic fidelity\" — reproduced real distributions",
        "Raises ethical questions about synthetic participants",
    ], subtitle_text="AI PERSONAS AS PARTICIPANTS")

    add_content_slide(prs, "Example 4: Predicting Human Decisions", [
        "Auletta, Kallen, di Bernardo, & Richardson (2023)",
        "LSTM neural networks + explainable AI (SHAP)",
        "Predicted expert vs novice decisions in collaborative task",
        "At timescales preceding conscious intent",
        "Experts more attuned to co-actor's behaviour —",
        "something traditional analysis missed",
    ], subtitle_text="ML FOR UNDERSTANDING BEHAVIOUR")

    add_content_slide(prs, "A Word of Caution", [
        "These are tools, not magic",
        "ML can find patterns in noise — including fake patterns",
        "Good research questions and critical thinking matter more than ever",
        "We'll spend as much time evaluating results as building models",
    ], subtitle_text="TOOLS, NOT MAGIC")

    # ============================================================
    # PART 8: MISCONCEPTIONS + WRAP-UP (~6 slides)
    # ============================================================

    add_section_slide(prs, "Common Misconceptions")

    add_two_column_text_slide(
        prs,
        "Four Things People Get Wrong",
        "Misconception",
        [
            "\"AI will replace researchers\"",
            "",
            "\"ML is just statistics\"",
            "",
            "\"You need to be a programmer\"",
            "",
            "\"AI is objective\"",
        ],
        "Reality",
        [
            "AI augments — handles computation",
            "so you focus on questions & design",
            "Overlap, but different emphases",
            "(inference vs prediction)",
            "Not with vibe coding & LLM assistants",
            "(but you must understand the output)",
            "Models inherit biases from training",
            "data and design choices",
        ],
    )

    add_content_slide(prs, "Getting Ready for Week 2", [
        "Complete the Getting Started setup guide (homework)",
        "GitHub account → Student Developer Pack → VS Code → Python",
        "Run the setup script",
        "Sign up for ChatGPT, Claude, or Gemini",
        "If stuck — bring your laptop to class",
    ], subtitle_text="BEFORE NEXT WEEK")

    add_content_slide(prs, "What's Ahead", [
        "Week 2: Setup troubleshooting + your first coding challenge",
        "Weeks 3–10: Regression, classification, trees, clustering, "
        "neural networks",
        "Week 11: Embeddings and LLMs — what's under the hood",
        "Weeks 12–13: Viva review and ethics discussion",
        "The LLM Problem-Solving Loop in every lab",
    ], subtitle_text="THE SEMESTER PLAN")

    add_section_slide(prs, "Questions?",
                      "Complete the setup guide before Week 2")

    # ============================================================
    # SAVE
    # ============================================================
    prs.save(OUTPUT_FILE)
    print(f"Slide deck saved to: {OUTPUT_FILE}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
